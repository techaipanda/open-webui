"""
RAG 文档加载器模块
功能: 从不同来源加载文档（PDF、网页、YouTube 等）
支持类型: PDF, URL, YouTube, GitHub, ArXiv, Office文档, EPUB 等

核心概念:
- 分块(Chunking): 将长文档分割成小块以便向量化和检索
- 向量化(Embedding): 将文本转换为数学向量表示
- 加载器工厂: 根据文件类型自动选择合适的加载器

依赖:
- langchain: 文档加载框架
- PyMuPDF: PDF 解析
- unstructured: 通用文档解析
"""

import asyncio
import requests
import logging
import ftfy
import sys
import json

from azure.identity import DefaultAzureCredential
from langchain_community.document_loaders import (
    AzureAIDocumentIntelligenceLoader,
    BSHTMLLoader,
    CSVLoader,
    Docx2txtLoader,
    OutlookMessageLoader,
    PyPDFLoader,
    TextLoader,
    YoutubeLoader,
)
from langchain_core.documents import Document

from open_webui.retrieval.loaders.external_document import ExternalDocumentLoader

from open_webui.retrieval.loaders.mistral import MistralLoader
from open_webui.retrieval.loaders.datalab_marker import DatalabMarkerLoader
from open_webui.retrieval.loaders.mineru import MinerULoader
from open_webui.retrieval.loaders.paddleocr_vl import PaddleOCRVLLoader

from open_webui.env import GLOBAL_LOG_LEVEL, REQUESTS_VERIFY, AIOHTTP_CLIENT_SESSION_SSL

logging.basicConfig(stream=sys.stdout, level=GLOBAL_LOG_LEVEL)
log = logging.getLogger(__name__)

# 支持的纯文本文件扩展名列表
# 用于区分需要特殊加载器的二进制文件
known_source_ext = [
    'go',
    'py',
    'java',
    'sh',
    'bat',
    'ps1',
    'cmd',
    'js',
    'ts',
    'css',
    'cpp',
    'hpp',
    'h',
    'c',
    'cs',
    'sql',
    'log',
    'ini',
    'pl',
    'pm',
    'r',
    'dart',
    'dockerfile',
    'env',
    'php',
    'hs',
    'hsc',
    'lua',
    'nginxconf',
    'conf',
    'm',
    'mm',
    'plsql',
    'perl',
    'rb',
    'rs',
    'db2',
    'scala',
    'bash',
    'swift',
    'vue',
    'svelte',
    'ex',
    'exs',
    'erl',
    'tsx',
    'jsx',
    'hs',
    'lhs',
    'json',
    'yaml',
    'yml',
    'toml',
]


class ExcelLoader:
    """
    Excel 文档加载器（备用方案）

    功能:
        使用 pandas 库从 Excel 文件中提取文本内容

    适用场景:
        当 unstructured 包未安装时的备用方案

    方法:
        load(): 读取 Excel 文件的每个工作表，返回合并的文本内容
    """

    def __init__(self, file_path):
        self.file_path = file_path

    def load(self) -> list[Document]:
        import pandas as pd

        text_parts = []
        xls = pd.ExcelFile(self.file_path)
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            text_parts.append(f'Sheet: {sheet_name}\n{df.to_string(index=False)}')
        return [
            Document(
                page_content='\n\n'.join(text_parts),
                metadata={'source': self.file_path},
            )
        ]


class PptxLoader:
    """
    PowerPoint 演示文稿加载器（备用方案）

    功能:
        使用 python-pptx 库从 PPT 文件中提取文本内容

    适用场景:
        当 unstructured 包未安装时的备用方案

    方法:
        load(): 读取演示文稿的每个幻灯片，返回合并的文本内容
    """

    def __init__(self, file_path):
        self.file_path = file_path

    def load(self) -> list[Document]:
        from pptx import Presentation

        prs = Presentation(self.file_path)
        text_parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_texts.append(shape.text_frame.text)
            if slide_texts:
                text_parts.append(f'Slide {i}:\n' + '\n'.join(slide_texts))
        return [
            Document(
                page_content='\n\n'.join(text_parts),
                metadata={'source': self.file_path},
            )
        ]


class TikaLoader:
    """
    Apache Tika 文档加载器

    功能:
        使用 Apache Tika 服务器从任意格式文档中提取文本

    支持格式:
        PDF, DOC, DOCX, XLS, XLSX, PPT, PPTX, EPUB, HTML, XML, JSON 等

    依赖:
        - Tika 服务器运行在指定 URL
        - Tika Python 客户端

    方法:
        load(): 通过 Tika REST API 提取文档文本
    """

    def __init__(self, url, file_path, mime_type=None, extract_images=None):
        self.url = url
        self.file_path = file_path
        self.mime_type = mime_type
        self.extract_images = extract_images

    def load(self) -> list[Document]:
        with open(self.file_path, 'rb') as f:
            data = f.read()

        if self.mime_type is not None:
            headers = {'Content-Type': self.mime_type}
        else:
            headers = {}

        if self.extract_images == True:
            headers['X-Tika-PDFextractInlineImages'] = 'true'

        endpoint = self.url
        if not endpoint.endswith('/'):
            endpoint += '/'
        endpoint += 'tika/text'

        r = requests.put(endpoint, data=data, headers=headers, verify=REQUESTS_VERIFY)

        if r.ok:
            raw_metadata = r.json()
            text = raw_metadata.get('X-TIKA:content', '<No text content found>').strip()

            if 'Content-Type' in raw_metadata:
                headers['Content-Type'] = raw_metadata['Content-Type']

            log.debug('Tika extracted text: %s', text)

            return [Document(page_content=text, metadata=headers)]
        else:
            raise Exception(f'Error calling Tika: {r.reason}')


class DoclingLoader:
    """
    Docling 文档智能解析加载器

    功能:
        使用 Docling API 将文档（PDF, DOCX,图片等）转换为 Markdown 格式

    特点:
        - 支持多种文档格式
        - 自动布局分析
        - 将表格转换为结构化文本

    方法:
        load(): 调用 Docling API 进行文档转换
    """

    def __init__(self, url, api_key=None, file_path=None, mime_type=None, params=None):
        self.url = url.rstrip('/')
        self.api_key = api_key
        self.file_path = file_path
        self.mime_type = mime_type

        self.params = params or {}

    def load(self) -> list[Document]:
        with open(self.file_path, 'rb') as f:
            headers = {}
            if self.api_key:
                headers['X-Api-Key'] = f'{self.api_key}'

            r = requests.post(
                f'{self.url}/v1/convert/file',
                files={
                    'files': (
                        self.file_path,
                        f,
                        self.mime_type or 'application/octet-stream',
                    )
                },
                data={
                    'image_export_mode': 'placeholder',
                    **self.params,
                },
                headers=headers,
                verify=AIOHTTP_CLIENT_SESSION_SSL,
            )
        if r.ok:
            result = r.json()
            document_data = result.get('document', {})
            text = document_data.get('md_content', '<No text content found>')

            metadata = {'Content-Type': self.mime_type} if self.mime_type else {}

            log.debug('Docling extracted text: %s', text)
            return [Document(page_content=text, metadata=metadata)]
        else:
            error_msg = f'Error calling Docling API: {r.reason}'
            if r.text:
                try:
                    error_data = r.json()
                    if 'detail' in error_data:
                        error_msg += f' - {error_data["detail"]}'
                except Exception:
                    error_msg += f' - {r.text}'
            raise Exception(f'Error calling Docling: {error_msg}')


class Loader:
    """
    文档加载器工厂类

    功能:
        根据文件类型和配置自动选择合适的加载器解析文档

    支持的加载引擎:
        - external: 外部文档加载服务
        - tika: Apache Tika 服务器
        - datalab_marker: DataLab Marker（支持 PDF/Office/EPUB 等）
        - docling: Docling 智能文档解析
        - document_intelligence: Azure AI Document Intelligence
        - mineru: MinerU 文档解析
        - mistral_ocr: Mistral OCR（PDF OCR）
        - paddleocr_vl: PaddleOCR 视觉语言模型

    核心逻辑:
        1. 根据文件扩展名和 MIME 类型判断文件类型
        2. 根据配置的引擎选择对应的加载器
        3. 调用加载器解析文档并返回 Document 对象列表

    方法:
        load(): 同步加载文档
        aload(): 异步加载文档（在线程池中执行，避免阻塞事件循环）
        _get_loader(): 根据文件类型获取合适的加载器实例
    """

    def __init__(self, engine: str = '', **kwargs):
        self.engine = engine
        self.user = kwargs.get('user', None)
        self.kwargs = kwargs

    def load(self, filename: str, file_content_type: str, file_path: str) -> list[Document]:
        """
        加载文档并返回 Document 对象列表

        Args:
            filename: 文件名
            file_content_type: 文件的 MIME 类型
            file_path: 文件路径

        Returns:
            Document 对象列表，每个 Document 包含 page_content 和 metadata
        """
        loader = self._get_loader(filename, file_content_type, file_path)
        docs = loader.load()

        return [Document(page_content=ftfy.fix_text(doc.page_content), metadata=doc.metadata) for doc in docs]

    async def aload(self, filename: str, file_content_type: str, file_path: str) -> list[Document]:
        """
        异步加载文档

        使用 asyncio.to_thread 在线程池中执行耗时的文档解析操作，
        避免阻塞异步事件循环

        Args:
            filename: 文件名
            file_content_type: 文件的 MIME 类型
            file_path: 文件路径

        Returns:
            Document 对象列表
        """
        return await asyncio.to_thread(self.load, filename, file_content_type, file_path)

    def _is_text_file(self, file_ext: str, file_content_type: str) -> bool:
        """
        判断是否为纯文本文件

        Args:
            file_ext: 文件扩展名
            file_content_type: MIME 类型

        Returns:
            是否为纯文本文件
        """
        return file_ext in known_source_ext or (
            file_content_type
            and file_content_type.find('text/') >= 0
            # Avoid text/html files being detected as text
            and not file_content_type.find('html') >= 0
        )

    def _get_loader(self, filename: str, file_content_type: str, file_path: str):
        file_ext = filename.split('.')[-1].lower()

        if (
            self.engine == 'external'
            and self.kwargs.get('EXTERNAL_DOCUMENT_LOADER_URL')
            and self.kwargs.get('EXTERNAL_DOCUMENT_LOADER_API_KEY')
        ):
            loader = ExternalDocumentLoader(
                file_path=file_path,
                url=self.kwargs.get('EXTERNAL_DOCUMENT_LOADER_URL'),
                api_key=self.kwargs.get('EXTERNAL_DOCUMENT_LOADER_API_KEY'),
                mime_type=file_content_type,
                user=self.user,
            )
        elif self.engine == 'tika' and self.kwargs.get('TIKA_SERVER_URL'):
            if self._is_text_file(file_ext, file_content_type):
                loader = TextLoader(file_path, autodetect_encoding=True)
            else:
                loader = TikaLoader(
                    url=self.kwargs.get('TIKA_SERVER_URL'),
                    file_path=file_path,
                    extract_images=self.kwargs.get('PDF_EXTRACT_IMAGES'),
                )
        elif (
            self.engine == 'datalab_marker'
            and self.kwargs.get('DATALAB_MARKER_API_KEY')
            and file_ext
            in [
                'pdf',
                'xls',
                'xlsx',
                'ods',
                'doc',
                'docx',
                'odt',
                'ppt',
                'pptx',
                'odp',
                'html',
                'epub',
                'png',
                'jpeg',
                'jpg',
                'webp',
                'gif',
                'tiff',
            ]
        ):
            api_base_url = self.kwargs.get('DATALAB_MARKER_API_BASE_URL', '')
            if not api_base_url or api_base_url.strip() == '':
                api_base_url = 'https://www.datalab.to/api/v1/marker'  # https://github.com/open-webui/open-webui/pull/16867#issuecomment-3218424349

            loader = DatalabMarkerLoader(
                file_path=file_path,
                api_key=self.kwargs['DATALAB_MARKER_API_KEY'],
                api_base_url=api_base_url,
                additional_config=self.kwargs.get('DATALAB_MARKER_ADDITIONAL_CONFIG'),
                use_llm=self.kwargs.get('DATALAB_MARKER_USE_LLM', False),
                skip_cache=self.kwargs.get('DATALAB_MARKER_SKIP_CACHE', False),
                force_ocr=self.kwargs.get('DATALAB_MARKER_FORCE_OCR', False),
                paginate=self.kwargs.get('DATALAB_MARKER_PAGINATE', False),
                strip_existing_ocr=self.kwargs.get('DATALAB_MARKER_STRIP_EXISTING_OCR', False),
                disable_image_extraction=self.kwargs.get('DATALAB_MARKER_DISABLE_IMAGE_EXTRACTION', False),
                format_lines=self.kwargs.get('DATALAB_MARKER_FORMAT_LINES', False),
                output_format=self.kwargs.get('DATALAB_MARKER_OUTPUT_FORMAT', 'markdown'),
            )
        elif self.engine == 'docling' and self.kwargs.get('DOCLING_SERVER_URL'):
            if self._is_text_file(file_ext, file_content_type):
                loader = TextLoader(file_path, autodetect_encoding=True)
            else:
                # Build params for DoclingLoader
                params = self.kwargs.get('DOCLING_PARAMS', {})
                if not isinstance(params, dict):
                    try:
                        params = json.loads(params)
                    except json.JSONDecodeError:
                        log.error('Invalid DOCLING_PARAMS format, expected JSON object')
                        params = {}

                loader = DoclingLoader(
                    url=self.kwargs.get('DOCLING_SERVER_URL'),
                    api_key=self.kwargs.get('DOCLING_API_KEY', None),
                    file_path=file_path,
                    mime_type=file_content_type,
                    params=params,
                )
        elif (
            self.engine == 'document_intelligence'
            and self.kwargs.get('DOCUMENT_INTELLIGENCE_ENDPOINT') != ''
            and (
                file_ext in ['pdf', 'docx', 'ppt', 'pptx']
                or file_content_type
                in [
                    'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                    'application/vnd.ms-powerpoint',
                    'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                ]
            )
        ):
            if self.kwargs.get('DOCUMENT_INTELLIGENCE_KEY') != '':
                loader = AzureAIDocumentIntelligenceLoader(
                    file_path=file_path,
                    api_endpoint=self.kwargs.get('DOCUMENT_INTELLIGENCE_ENDPOINT'),
                    api_key=self.kwargs.get('DOCUMENT_INTELLIGENCE_KEY'),
                    api_model=self.kwargs.get('DOCUMENT_INTELLIGENCE_MODEL'),
                )
            else:
                loader = AzureAIDocumentIntelligenceLoader(
                    file_path=file_path,
                    api_endpoint=self.kwargs.get('DOCUMENT_INTELLIGENCE_ENDPOINT'),
                    azure_credential=DefaultAzureCredential(),
                    api_model=self.kwargs.get('DOCUMENT_INTELLIGENCE_MODEL'),
                )
        elif self.engine == 'mineru' and file_ext in ['pdf']:  # MinerU currently only supports PDF
            mineru_timeout = self.kwargs.get('MINERU_API_TIMEOUT', 300)
            if mineru_timeout:
                try:
                    mineru_timeout = int(mineru_timeout)
                except ValueError:
                    mineru_timeout = 300

            loader = MinerULoader(
                file_path=file_path,
                api_mode=self.kwargs.get('MINERU_API_MODE', 'local'),
                api_url=self.kwargs.get('MINERU_API_URL', 'http://localhost:8000'),
                api_key=self.kwargs.get('MINERU_API_KEY', ''),
                params=self.kwargs.get('MINERU_PARAMS', {}),
                timeout=mineru_timeout,
            )
        elif (
            self.engine == 'mistral_ocr'
            and self.kwargs.get('MISTRAL_OCR_API_KEY') != ''
            and file_ext in ['pdf']  # Mistral OCR currently only supports PDF and images
        ):
            loader = MistralLoader(
                base_url=self.kwargs.get('MISTRAL_OCR_API_BASE_URL'),
                api_key=self.kwargs.get('MISTRAL_OCR_API_KEY'),
                file_path=file_path,
            )
        elif self.engine == 'paddleocr_vl' and self.kwargs.get('PADDLEOCR_VL_TOKEN') != '':
            loader = PaddleOCRVLLoader(
                api_url=self.kwargs.get('PADDLEOCR_VL_BASE_URL'),
                token=self.kwargs.get('PADDLEOCR_VL_TOKEN'),
                file_path=file_path,
            )
        else:
            if file_ext == 'pdf':
                loader = PyPDFLoader(
                    file_path,
                    extract_images=self.kwargs.get('PDF_EXTRACT_IMAGES'),
                    mode=self.kwargs.get('PDF_LOADER_MODE', 'page'),
                )
            elif file_ext == 'csv':
                loader = CSVLoader(file_path, autodetect_encoding=True)
            elif file_ext == 'rst':
                try:
                    from langchain_community.document_loaders import UnstructuredRSTLoader

                    loader = UnstructuredRSTLoader(file_path, mode='elements')
                except ImportError:
                    log.warning(
                        "The 'unstructured' package is not installed. "
                        'Falling back to plain text loading for .rst file. '
                        'Install it with: pip install unstructured'
                    )
                    loader = TextLoader(file_path, autodetect_encoding=True)
            elif file_ext == 'xml':
                try:
                    from langchain_community.document_loaders import UnstructuredXMLLoader

                    loader = UnstructuredXMLLoader(file_path)
                except ImportError:
                    log.warning(
                        "The 'unstructured' package is not installed. "
                        'Falling back to plain text loading for .xml file. '
                        'Install it with: pip install unstructured'
                    )
                    loader = TextLoader(file_path, autodetect_encoding=True)
            elif file_ext in ['htm', 'html']:
                loader = BSHTMLLoader(file_path, open_encoding='unicode_escape')
            elif file_ext == 'md':
                loader = TextLoader(file_path, autodetect_encoding=True)
            elif file_content_type == 'application/epub+zip':
                try:
                    from langchain_community.document_loaders import UnstructuredEPubLoader

                    loader = UnstructuredEPubLoader(file_path)
                except ImportError:
                    raise ValueError(
                        "Processing .epub files requires the 'unstructured' package. "
                        'Install it with: pip install unstructured'
                    )
            elif (
                file_content_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                or file_ext == 'docx'
            ):
                loader = Docx2txtLoader(file_path)
            elif file_content_type in [
                'application/vnd.ms-excel',
                'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            ] or file_ext in ['xls', 'xlsx']:
                try:
                    from langchain_community.document_loaders import UnstructuredExcelLoader

                    loader = UnstructuredExcelLoader(file_path)
                except ImportError:
                    log.warning(
                        "The 'unstructured' package is not installed. "
                        'Falling back to pandas for Excel file loading. '
                        'Install unstructured for better results: pip install unstructured'
                    )
                    loader = ExcelLoader(file_path)
            elif file_content_type in [
                'application/vnd.ms-powerpoint',
                'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            ] or file_ext in ['ppt', 'pptx']:
                try:
                    from langchain_community.document_loaders import UnstructuredPowerPointLoader

                    loader = UnstructuredPowerPointLoader(file_path)
                except ImportError:
                    log.warning(
                        "The 'unstructured' package is not installed. "
                        'Falling back to python-pptx for PowerPoint file loading. '
                        'Install unstructured for better results: pip install unstructured'
                    )
                    loader = PptxLoader(file_path)
            elif file_ext == 'msg':
                loader = OutlookMessageLoader(file_path)
            elif file_ext == 'odt':
                try:
                    from langchain_community.document_loaders import UnstructuredODTLoader

                    loader = UnstructuredODTLoader(file_path)
                except ImportError:
                    raise ValueError(
                        "Processing .odt files requires the 'unstructured' package. "
                        'Install it with: pip install unstructured'
                    )
            elif self._is_text_file(file_ext, file_content_type):
                loader = TextLoader(file_path, autodetect_encoding=True)
            else:
                loader = TextLoader(file_path, autodetect_encoding=True)

        return loader
